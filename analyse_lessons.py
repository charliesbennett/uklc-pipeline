"""
analyse_lessons.py
------------------
Run this once after adding your PPTX files to the knowledge/ folder.
It extracts structure, content patterns and style rules from all real lessons,
then updates the generator agent's system prompt automatically.

Usage:
    cd ~/Documents/uklc-pipeline
    python3 analyse_lessons.py
"""

import json
import os
import re
import sys
from pathlib import Path
from collections import Counter, defaultdict

# Try markitdown for text extraction
try:
    from markitdown import MarkItDown
    md = MarkItDown()
    USE_MARKITDOWN = True
except ImportError:
    USE_MARKITDOWN = False

try:
    from pptx import Presentation
    USE_PPTX = True
except ImportError:
    USE_PPTX = False

KNOWLEDGE_DIR = Path(__file__).parent / "knowledge"
PPTX_FILES    = list(KNOWLEDGE_DIR.glob("*.pptx")) + list(KNOWLEDGE_DIR.glob("*.PPTX"))
OUTPUT_PATH   = KNOWLEDGE_DIR / "style_guide.json"
AGENT_PATH    = Path(__file__).parent / "agents" / "generator.py"


def extract_slides(pptx_path: Path) -> list[dict]:
    """Extract all slide text from a PPTX file."""
    slides = []
    if not USE_PPTX:
        return slides
    try:
        prs = Presentation(str(pptx_path))
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in para.runs).strip()
                        if line:
                            texts.append(line)
            slides.append({"index": i, "texts": texts, "full": "\n".join(texts)})
    except Exception as e:
        print(f"  Warning: could not read {pptx_path.name}: {e}")
    return slides


def classify_slide(slide: dict) -> str:
    """Classify a slide by its content."""
    full = slide["full"].lower()
    texts = slide["texts"]
    first = texts[0].lower() if texts else ""

    if not texts:
        return "blank"
    if any(x in full for x in ["lesson summary", "lesson focus", "objectives", "extra materials"]):
        return "summary"
    if "teacher's notes" in full or "feel free to use this space" in full:
        return "teacher_notes"
    if "notes on materials" in full or "materials" == first:
        return "materials"
    if "this slide is meant to be hidden" in full:
        return "hidden_plan"
    if re.search(r'\d+\.\s+(introduction|vocabulary|reading|task|feedback|discussion|analysis)', full):
        return "hidden_plan"
    if any(x in full for x in ["week a", "week b", "week c", "week d", "purpose | lang", "purpose | lead", "purpose | cult"]):
        return "title"
    if any(x in full for x in ["discuss in pairs", "discuss in groups", "what would you", "have you ever", "what do you think"]):
        if len(texts) <= 5:
            return "warmup"
        return "hook"
    if "watch" in full and any(x in full for x in ["spot", "note down", "as you watch"]):
        return "video_task"
    if re.search(r'^\s*video\s*$', full, re.MULTILINE) or "embed" in full or "youtube" in full.lower():
        return "video_placeholder"
    if "answers:" in full or (first == "answers" and len(texts) < 12):
        return "answers_hidden"
    if any(x in full for x in ["fill in", "complete the sentence", "use the word", "gap"]):
        return "gap_fill"
    if any(x in full for x in ["match", "matching"]):
        return "matching"
    if any(x in full for x in ["quiz", "buzzer", "round 1", "round 2", "q1:", "q2:"]):
        return "quiz"
    if re.search(r'^[1-5]\s', "\n".join(texts[:6]), re.MULTILINE) and len(texts) >= 4:
        return "discussion"
    if any(x in full for x in ["vocabulary", "key words", "language for", "word bank"]):
        return "vocab_intro"
    if any(x in full for x in ["in 20", "in 201", "in 202", "billion", "million", "per cent", "%", "according to"]):
        return "info"
    if any(x in full for x in ["your task", "your challenge", "in your group", "you are going to", "your group will"]):
        return "task_setup"
    if any(x in full for x in ["use this structure", "step 1", "step 2", "role", "presenter", "demonstrator"]):
        return "task_content"
    if any(x in full for x in ["vote", "feedback", "award", "star:", "suggestion:"]):
        return "feedback"
    if any(x in full for x in ["reflection", "complete one sentence", "what did you learn"]):
        return "reflection"
    if any(x in full for x in ["debate", "group a", "group b", "argue"]):
        return "debate"
    return "content"


def analyse_lesson(pptx_path: Path) -> dict:
    """Analyse one PPTX and return its pattern data."""
    slides = extract_slides(pptx_path)
    if not slides:
        return {}

    classified = [(s, classify_slide(s)) for s in slides]
    student_slides = [(s, t) for s, t in classified
                      if t not in ("title", "summary", "hidden_plan", "materials", "teacher_notes")]

    # Extract slide sequence
    sequence = [t for _, t in student_slides]

    # Extract discussion questions
    discussion_questions = []
    for s, t in classified:
        if t == "discussion":
            qs = [l for l in s["texts"] if re.match(r'^\d+[\.\)]?\s+\w', l) or
                  (l.endswith("?") and len(l) > 20)]
            discussion_questions.extend(qs[:6])

    # Extract vocab items (word — definition pattern)
    vocab_items = []
    for s, t in classified:
        if t == "vocab_intro":
            for line in s["texts"]:
                if " — " in line or " - " in line or ": " in line:
                    vocab_items.append(line[:100])

    # Extract gap fill sentences
    gap_fill_sentences = []
    for s, t in classified:
        if t == "gap_fill":
            for line in s["texts"]:
                if "___" in line or "___ " in line or re.search(r'\b___+\b', line):
                    gap_fill_sentences.append(line[:120])

    # Extract teacher instructions style
    teacher_instructions = []
    for s, t in classified:
        if t == "hidden_plan":
            for line in s["texts"]:
                if len(line) > 30 and not any(x in line.lower() for x in
                    ["time", "shortcut", "this slide", "hidden"]):
                    teacher_instructions.append(line[:200])

    # Check for key patterns
    has_video    = "video_placeholder" in sequence or "video_task" in sequence
    has_answers  = "answers_hidden" in sequence
    has_warmup   = "warmup" in sequence
    ends_with    = sequence[-1] if sequence else "unknown"
    has_reflect  = "reflection" in sequence

    return {
        "filename": pptx_path.name,
        "total_slides": len(slides),
        "student_slide_count": len(student_slides),
        "sequence": sequence,
        "has_video_placeholder": has_video,
        "has_hidden_answers": has_answers,
        "has_warmup": has_warmup,
        "ends_with": ends_with,
        "has_reflection": has_reflect,
        "sample_discussion_questions": discussion_questions[:5],
        "sample_vocab_items": vocab_items[:6],
        "sample_gap_fills": gap_fill_sentences[:4],
        "sample_teacher_instructions": teacher_instructions[:3],
    }


def build_style_guide(analyses: list[dict]) -> dict:
    """Aggregate patterns from all lessons into a style guide."""
    if not analyses:
        return {}

    # Slide type frequency
    all_sequences = [a["sequence"] for a in analyses if a.get("sequence")]
    type_counter  = Counter(t for seq in all_sequences for t in seq)
    slide_counts  = [a["student_slide_count"] for a in analyses if a.get("student_slide_count")]

    # Most common sequence patterns
    # Find most common slide types in order position
    max_len = max((len(s) for s in all_sequences), default=0)
    position_types = defaultdict(list)
    for seq in all_sequences:
        for i, t in enumerate(seq):
            position_types[i].append(t)
    typical_sequence = [Counter(position_types[i]).most_common(1)[0][0]
                        for i in range(min(max_len, 18))]

    # Pattern statistics
    has_video_pct   = sum(1 for a in analyses if a.get("has_video_placeholder")) / len(analyses) * 100
    has_answers_pct = sum(1 for a in analyses if a.get("has_hidden_answers"))    / len(analyses) * 100
    has_warmup_pct  = sum(1 for a in analyses if a.get("has_warmup"))            / len(analyses) * 100
    has_reflect_pct = sum(1 for a in analyses if a.get("has_reflection"))        / len(analyses) * 100

    ending_types = Counter(a.get("ends_with") for a in analyses if a.get("ends_with"))

    # Collect all examples
    all_discussion_qs   = [q for a in analyses for q in a.get("sample_discussion_questions", [])]
    all_vocab           = [v for a in analyses for v in a.get("sample_vocab_items", [])]
    all_gap_fills       = [g for a in analyses for g in a.get("sample_gap_fills", [])]
    all_teacher_instrs  = [t for a in analyses for t in a.get("sample_teacher_instructions", [])]

    return {
        "lessons_analysed": len(analyses),
        "filenames": [a["filename"] for a in analyses],
        "slide_counts": {
            "min": min(slide_counts) if slide_counts else 0,
            "max": max(slide_counts) if slide_counts else 0,
            "avg": round(sum(slide_counts) / len(slide_counts), 1) if slide_counts else 0,
        },
        "slide_type_frequency": dict(type_counter.most_common(20)),
        "typical_sequence": typical_sequence,
        "pattern_stats": {
            "has_video_placeholder_pct":  round(has_video_pct),
            "has_hidden_answers_pct":     round(has_answers_pct),
            "has_warmup_slide_pct":       round(has_warmup_pct),
            "has_reflection_pct":         round(has_reflect_pct),
            "most_common_ending":         ending_types.most_common(1)[0][0] if ending_types else "unknown",
        },
        "example_discussion_questions": all_discussion_qs[:15],
        "example_vocab_items":          all_vocab[:15],
        "example_gap_fill_sentences":   all_gap_fills[:10],
        "example_teacher_instructions": all_teacher_instrs[:8],
    }


def update_generator_prompt(style_guide: dict):
    """Inject the style guide summary into the generator agent."""
    if not AGENT_PATH.exists():
        print("  Warning: generator.py not found, skipping prompt update")
        return

    with open(AGENT_PATH) as f:
        code = f.read()

    summary = f"""
# ── LEARNT FROM {style_guide['lessons_analysed']} REAL LESSONS ─────────────────────────────────────────
# Slide counts: avg {style_guide['slide_counts']['avg']}, min {style_guide['slide_counts']['min']}, max {style_guide['slide_counts']['max']}
# Typical sequence: {' → '.join(style_guide['typical_sequence'][:16])}
# Video placeholder present: {style_guide['pattern_stats']['has_video_placeholder_pct']}% of lessons
# Hidden answers present:    {style_guide['pattern_stats']['has_hidden_answers_pct']}% of lessons
# Warm-up slide present:     {style_guide['pattern_stats']['has_warmup_slide_pct']}% of lessons
# Reflection at end:         {style_guide['pattern_stats']['has_reflection_pct']}% of lessons
# Most common ending slide:  {style_guide['pattern_stats']['most_common_ending']}
"""

    # Replace or insert the learnt block
    marker_start = "# ── LEARNT FROM"
    marker_end   = "# ─────────────────────────────────────────────────────────────────────\n"
    if marker_start in code:
        start = code.index(marker_start)
        end   = code.index(marker_end, start) + len(marker_end)
        code  = code[:start] + summary.strip() + "\n" + code[end:]
    else:
        # Insert after the imports block
        code = code.replace('SYSTEM_PROMPT = """', summary + '\nSYSTEM_PROMPT = """', 1)

    with open(AGENT_PATH, "w") as f:
        f.write(code)
    print("  ✓ Generator agent updated with style patterns")


def main():
    print(f"\n UKLC Purpose Lesson Style Analyser")
    print(f"{'─'*50}")

    if not PPTX_FILES:
        print(f"\n No PPTX files found in {KNOWLEDGE_DIR}")
        print("  → Add your lesson PPTXs to the knowledge/ folder and run again.")
        sys.exit(1)

    print(f"\n Found {len(PPTX_FILES)} PPTX file(s) in knowledge/\n")

    analyses = []
    for pptx in sorted(PPTX_FILES):
        print(f"  Analysing {pptx.name}…")
        result = analyse_lesson(pptx)
        if result:
            analyses.append(result)
            seq_preview = " → ".join(result["sequence"][:8])
            print(f"    {result['student_slide_count']} student slides: {seq_preview}…")

    if not analyses:
        print("\n  No lessons could be analysed. Check your PPTX files.")
        sys.exit(1)

    print(f"\n Building style guide from {len(analyses)} lessons…")
    style_guide = build_style_guide(analyses)

    # Save style guide
    with open(OUTPUT_PATH, "w") as f:
        json.dump(style_guide, f, indent=2)
    print(f"  ✓ Style guide saved to knowledge/style_guide.json")

    # Update generator
    print("  Updating generator agent…")
    update_generator_prompt(style_guide)

    # Print summary
    print(f"""
{'─'*50}
 RESULTS

  Lessons analysed:   {style_guide['lessons_analysed']}
  Avg student slides: {style_guide['slide_counts']['avg']}
  Typical sequence:   {' → '.join(style_guide['typical_sequence'][:10])}…

  Patterns found:
    Video placeholder: {style_guide['pattern_stats']['has_video_placeholder_pct']}% of lessons
    Hidden answers:    {style_guide['pattern_stats']['has_hidden_answers_pct']}% of lessons
    Warm-up slide:     {style_guide['pattern_stats']['has_warmup_slide_pct']}% of lessons
    Reflection at end: {style_guide['pattern_stats']['has_reflection_pct']}% of lessons
    Ends with:         {style_guide['pattern_stats']['most_common_ending']}

  Top slide types:
{''.join(f"    {t}: {n}x{chr(10)}" for t, n in list(style_guide['slide_type_frequency'].items())[:8])}
  The generator agent has been updated.
  Run python3 app.py to start generating with your patterns.
{'─'*50}
""")


if __name__ == "__main__":
    main()
